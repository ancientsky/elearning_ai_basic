"""shared.py — Brand constants, helpers, and reusable slide builders.

All module scripts import from this file.
Run with:  cd slides && uv run python generate_all.py
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand palette (from css/style.css :root) ──────────────────────────────────
PRIMARY   = RGBColor(0x1A, 0x5F, 0x7A)   # #1a5f7a
SECONDARY = RGBColor(0x57, 0xC5, 0xB6)   # #57c5b6
ACCENT    = RGBColor(0x15, 0x98, 0x95)   # #159895
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF4, 0xF8, 0xFA)
DARK_BG   = RGBColor(0x0D, 0x2B, 0x3E)   # deep navy for covers
MUTED     = RGBColor(0x6C, 0x75, 0x7D)
SUCCESS   = RGBColor(0x28, 0xA7, 0x45)
WARNING   = RGBColor(0xFF, 0xC1, 0x07)
DANGER    = RGBColor(0xDC, 0x35, 0x45)
CODE_BG   = RGBColor(0x2D, 0x2D, 0x2D)
CODE_FG   = RGBColor(0xF8, 0xF8, 0xF2)
ROW_ALT   = RGBColor(0xE4, 0xF3, 0xF8)

# RTFC element colours
ROLE_C   = RGBColor(0xE7, 0x4C, 0x3C)
TASK_C   = RGBColor(0x34, 0x98, 0xDB)
FORMAT_C = RGBColor(0x2E, 0xCC, 0x71)
CONST_C  = RGBColor(0x9B, 0x59, 0xB6)

# ── Slide dimensions (16:9) ────────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "output")


# ── Low-level helpers ──────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _apply_font(run, size: float, bold: bool, italic: bool,
                color: RGBColor | None) -> None:
    """Set font face (Microsoft JhengHei) + size/style via XML for correct CJK rendering."""
    run.font.name   = "Microsoft JhengHei"
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is not None:
            rPr.remove(el)
    for tag, tf in (("a:latin", "Microsoft JhengHei"),
                    ("a:ea",    "Microsoft JhengHei"),
                    ("a:cs",    "Microsoft JhengHei")):
        el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", tf)


def txbox(slide, text: str,
          left, top, width, height,
          size: float = 18,
          bold: bool = False,
          italic: bool = False,
          color: RGBColor = WHITE,
          align=PP_ALIGN.LEFT,
          wrap: bool = True):
    """Add a text-box and return the shape."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _apply_font(run, size, bold, italic, color)
    return tb


def rect(slide, left, top, width, height,
         fill: RGBColor, line: RGBColor | None = None):
    """Add a solid rectangle and return the shape."""
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp


def header_bar(slide, title: str, bar_color=PRIMARY,
               text_color=WHITE, font_size=26) -> None:
    """Full-width top bar with title text."""
    rect(slide, Inches(0), Inches(0), SW, Inches(1.15), bar_color)
    txbox(slide, title,
          Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.75),
          size=font_size, bold=True, color=text_color)


def save_pptx(prs: Presentation, filename: str) -> str:
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    path = os.path.join(OUTPUT_DIR, filename)
    prs.save(path)
    n = len(prs.slides)
    print(f"  ✅  {filename}  ({n} slides)")
    return path


# ── Reusable slide builders ────────────────────────────────────────────────────

def cover_slide(prs, module_num: str, title: str,
                subtitle: str, duration: str):
    """Dark-navy cover with left teal accent bar."""
    sl = blank(prs)
    set_bg(sl, DARK_BG)

    # Left accent stripe
    rect(sl, Inches(0), Inches(0), Inches(0.22), SH, SECONDARY)

    # Module badge pill
    rect(sl, Inches(0.55), Inches(1.5), Inches(1.8), Inches(0.6), SECONDARY)
    txbox(sl, f"模組 {module_num}",
          Inches(0.55), Inches(1.5), Inches(1.8), Inches(0.6),
          size=18, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

    # Title
    txbox(sl, title,
          Inches(0.55), Inches(2.3), Inches(11.5), Inches(2.0),
          size=44, bold=True, color=WHITE)

    # Subtitle
    txbox(sl, subtitle,
          Inches(0.55), Inches(4.5), Inches(11.5), Inches(0.8),
          size=22, color=SECONDARY)

    # Duration
    txbox(sl, f"⏱  片長：{duration}",
          Inches(0.55), Inches(5.4), Inches(5.0), Inches(0.5),
          size=16, color=MUTED)

    # Footer
    txbox(sl, "衛生福利部疾病管制署  AI 應用基礎班",
          Inches(0.55), Inches(6.8), Inches(12.2), Inches(0.4),
          size=13, color=MUTED)
    return sl


def objectives_slide(prs, points: list[str],
                     header="完成後，您將能夠："):
    """Three learning-objective bullet slide."""
    sl = blank(prs)
    set_bg(sl, LIGHT_BG)
    header_bar(sl, header)

    colors = [ROLE_C, TASK_C, FORMAT_C]
    for i, pt in enumerate(points[:3]):
        top = Inches(1.5) + i * Inches(1.85)
        rect(sl, Inches(0.5), top, Inches(0.75), Inches(0.75), colors[i])
        txbox(sl, str(i + 1),
              Inches(0.5), top, Inches(0.75), Inches(0.75),
              size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(sl, pt,
              Inches(1.45), top + Inches(0.05),
              Inches(11.2), Inches(1.5),
              size=20, color=PRIMARY)
    return sl


def content_slide(prs, title: str, bullets: list[str]):
    """Standard bullet-list slide with header bar."""
    sl = blank(prs)
    set_bg(sl, LIGHT_BG)
    header_bar(sl, title)

    top = Inches(1.4)
    for b in bullets:
        indent = b.startswith("  ")
        text   = ("    " if indent else "• ") + b.lstrip()
        sz     = 17 if indent else 20
        clr    = MUTED if indent else PRIMARY
        txbox(sl, text,
              Inches(0.6), top, Inches(12.1), Inches(0.75),
              size=sz, color=clr)
        top += Inches(0.68)
    return sl


def two_col_slide(prs, title: str,
                  left_head: str, left_items: list[str],
                  right_head: str, right_items: list[str],
                  left_color=SUCCESS, right_color=DANGER):
    """Split ✓ / ✗ comparison slide."""
    sl = blank(prs)
    set_bg(sl, LIGHT_BG)
    header_bar(sl, title)

    for col_idx, (head, items, clr, icon, ox) in enumerate([
        (left_head,  left_items,  left_color,  "✓", Inches(0.4)),
        (right_head, right_items, right_color, "✗", Inches(7.0)),
    ]):
        rect(sl, ox, Inches(1.3), Inches(5.8), Inches(0.65), clr)
        txbox(sl, f"{icon}  {head}",
              ox, Inches(1.3), Inches(5.8), Inches(0.65),
              size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        top = Inches(2.1)
        for item in items:
            txbox(sl, f"{icon}  {item}",
                  ox + Inches(0.15), top,
                  Inches(5.5), Inches(0.7),
                  size=18, color=PRIMARY)
            top += Inches(0.7)
    return sl


def table_slide(prs, title: str,
                headers: list[str], rows: list[list[str]],
                col_widths: list | None = None):
    """Data table slide."""
    sl = blank(prs)
    set_bg(sl, LIGHT_BG)
    header_bar(sl, title)

    ncols = len(headers)
    nrows = len(rows)
    if col_widths is None:
        cw = Inches(12.3) / ncols
        col_widths = [cw] * ncols

    row_h = Inches(0.54)
    tbl_h = row_h * (nrows + 1)
    tbl = sl.shapes.add_table(
        nrows + 1, ncols,
        Inches(0.5), Inches(1.35),
        int(sum(col_widths)), int(tbl_h)
    ).table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = int(cw)

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        _apply_font(run, 16, True, False, WHITE)

    for ri, row in enumerate(rows):
        bg = LIGHT_BG if ri % 2 == 0 else ROW_ALT
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = val
            _apply_font(run, 15, False, False, PRIMARY)
    return sl


def code_slide(prs, title: str, code: str, caption: str = ""):
    """Dark code-box prompt example slide."""
    sl = blank(prs)
    set_bg(sl, LIGHT_BG)
    header_bar(sl, title)

    box_top = Inches(1.35)
    box_h   = Inches(5.55) if not caption else Inches(4.9)
    rect(sl, Inches(0.5), box_top, Inches(12.3), box_h, CODE_BG)
    txbox(sl, code,
          Inches(0.8), box_top + Inches(0.18),
          Inches(11.8), box_h - Inches(0.35),
          size=15, color=CODE_FG)
    if caption:
        txbox(sl, caption,
              Inches(0.5), Inches(6.5),
              Inches(12.3), Inches(0.65),
              size=16, italic=True, color=MUTED)
    return sl


def rtfc_card_slide(prs):
    """Four-colour RTFC overview card."""
    sl = blank(prs)
    set_bg(sl, LIGHT_BG)
    header_bar(sl, "Prompt 四要素：RTFC 框架")

    cards = [
        (ROLE_C,   "R", "Role  角色",
         "告訴 AI 扮演什麼身分\n例：你是疾管署衛教專家"),
        (TASK_C,   "T", "Task  任務",
         "具體說明要做什麼事\n例：請撰寫登革熱衛教文宣"),
        (FORMAT_C, "F", "Format  格式",
         "指定輸出的呈現方式\n例：分三段、使用條列式"),
        (CONST_C,  "C", "Constraint  限制",
         "設定邊界條件\n例：300 字以內，語氣親切"),
    ]
    for i, (clr, letter, name, desc) in enumerate(cards):
        col = i % 2
        row = i // 2
        left = Inches(0.45) + col * Inches(6.45)
        top  = Inches(1.4)  + row * Inches(2.85)
        rect(sl, left, top, Inches(6.1), Inches(2.6), clr)
        txbox(sl, letter,
              left + Inches(0.2), top + Inches(0.25),
              Inches(1.1), Inches(1.1),
              size=52, bold=True, color=WHITE)
        txbox(sl, name,
              left + Inches(1.35), top + Inches(0.25),
              Inches(4.5), Inches(0.65),
              size=22, bold=True, color=WHITE)
        txbox(sl, desc,
              left + Inches(1.35), top + Inches(1.0),
              Inches(4.5), Inches(1.4),
              size=16, color=WHITE)
    return sl


def cta_slide(prs, message: str, button_text: str):
    """Closing call-to-action slide."""
    sl = blank(prs)
    set_bg(sl, PRIMARY)

    txbox(sl, message,
          Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.0),
          size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rect(sl, Inches(3.3), Inches(4.1), Inches(6.7), Inches(0.9), SECONDARY)
    txbox(sl, button_text,
          Inches(3.3), Inches(4.1), Inches(6.7), Inches(0.9),
          size=22, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

    txbox(sl, "疾病管制署  AI 應用基礎班",
          Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
          size=13, color=SECONDARY, align=PP_ALIGN.CENTER)
    return sl


def alert_slide(prs, title: str, alert_type: str,
                content: str, icon: str = ""):
    """Full-bleed coloured alert slide."""
    sl = blank(prs)
    set_bg(sl, LIGHT_BG)
    header_bar(sl, title)

    colors = {
        "danger":  DANGER,
        "warning": WARNING,
        "success": SUCCESS,
        "info":    ACCENT,
    }
    aclr = colors.get(alert_type, ACCENT)
    rect(sl, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.0), aclr)
    prefix = (icon + "  ") if icon else ""
    txbox(sl, prefix + content,
          Inches(0.85), Inches(1.75),
          Inches(11.7), Inches(4.5),
          size=21, color=WHITE)
    return sl


def section_divider(prs, title: str, subtitle: str = ""):
    """Teal full-bleed section divider."""
    sl = blank(prs)
    set_bg(sl, SECONDARY)
    txbox(sl, title,
          Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.4),
          size=40, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    if subtitle:
        txbox(sl, subtitle,
              Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.8),
              size=22, color=DARK_BG, align=PP_ALIGN.CENTER)
    return sl


def card_grid_slide(prs, title: str,
                    cards: list[tuple[RGBColor, str, str]]):
    """Grid of coloured cards (up to 4).  card = (color, title, body)"""
    sl = blank(prs)
    set_bg(sl, LIGHT_BG)
    header_bar(sl, title)

    positions = [
        (Inches(0.45), Inches(1.4)),
        (Inches(7.0),  Inches(1.4)),
        (Inches(0.45), Inches(4.3)),
        (Inches(7.0),  Inches(4.3)),
    ]
    for i, (clr, ctitle, cbody) in enumerate(cards[:4]):
        left, top = positions[i]
        rect(sl, left, top, Inches(5.9), Inches(2.7), clr)
        txbox(sl, ctitle,
              left + Inches(0.25), top + Inches(0.25),
              Inches(5.4), Inches(0.7),
              size=20, bold=True, color=WHITE)
        txbox(sl, cbody,
              left + Inches(0.25), top + Inches(1.0),
              Inches(5.4), Inches(1.5),
              size=16, color=WHITE)
    return sl


def dialogue_step_slide(prs, step: int, total: int,
                        round_label: str, user_text: str,
                        ai_note: str = ""):
    """Multi-turn dialogue step slide."""
    sl = blank(prs)
    set_bg(sl, LIGHT_BG)
    header_bar(sl, f"多輪對話範例  （第 {step}/{total} 輪）")

    # Step indicator dots
    for j in range(1, total + 1):
        clr = SECONDARY if j == step else RGBColor(0xCC, 0xCC, 0xCC)
        rect(sl, Inches(0.5) + (j - 1) * Inches(0.55),
             Inches(1.25), Inches(0.42), Inches(0.42), clr)
        txbox(sl, str(j),
              Inches(0.5) + (j - 1) * Inches(0.55),
              Inches(1.25), Inches(0.42), Inches(0.42),
              size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Round label
    txbox(sl, round_label,
          Inches(0.5), Inches(1.85), Inches(12.0), Inches(0.55),
          size=18, bold=True, color=PRIMARY)

    # User prompt box (dark)
    rect(sl, Inches(0.5), Inches(2.5), Inches(12.3), Inches(2.5), CODE_BG)
    txbox(sl, user_text,
          Inches(0.8), Inches(2.65), Inches(11.8), Inches(2.2),
          size=16, color=CODE_FG)

    # AI note
    if ai_note:
        rect(sl, Inches(0.5), Inches(5.15), Inches(12.3), Inches(1.3), SECONDARY)
        txbox(sl, "AI 觀察：" + ai_note,
              Inches(0.75), Inches(5.25), Inches(11.8), Inches(1.15),
              size=17, color=DARK_BG)
    return sl
